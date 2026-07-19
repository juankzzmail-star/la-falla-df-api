"""Tests for the shared multi-format extractor (change multi-format-ingestion).

Authored in the `api_gateway` package layout (CI/container). No network and no production Postgres:
the vision model (OpenAI) and the link fetcher (Firecrawl/requests) are monkeypatched; heavy parsers
(openpyxl, python-pptx) are `importorskip`-ed so this file runs in the built image and skips locally.

Covers tasks.md §4.2/§4.3/§4.4:
- extract per format: XLSX -> table per sheet, PPTX -> slides+notes, CSV -> table, image -> vision text
  (mocked), URL -> markdown (firecrawl mocked) and failure -> "", unsupported binary -> "" (honest)
- DB state: an xlsx and a (mocked) link end up as chunks in document_chunks via ingest-resource
- endpoint: /strategy/ingest-resource (obs with xlsx; url) + auth rejection without key
"""
import io
import os

os.environ.setdefault("DATABASE_URL", "sqlite://")
os.environ.setdefault("FASTAPI_GM_API_KEY", "test-key-123")
os.environ.setdefault("FASTAPI_API_KEY", "test-key-123")
os.environ.setdefault("OPENAI_API_KEY", "test-openai")  # presence only; vision/embeddings are mocked

import pytest
from sqlalchemy import create_engine, text

from api_gateway.routers import _extract
from api_gateway.routers import rag as rag_mod
from api_gateway.routers import strategy as st

H = {"X-API-Key": "test-key-123"}

_DDL_CHUNKS = """
CREATE TABLE document_chunks (
  id INTEGER PRIMARY KEY AUTOINCREMENT,
  source_type TEXT NOT NULL,
  source_id   TEXT,
  source_name TEXT NOT NULL,
  chunk_index INTEGER NOT NULL,
  chunk_text  TEXT NOT NULL,
  embedding   TEXT,
  metadata    TEXT,
  created_at  TIMESTAMP DEFAULT CURRENT_TIMESTAMP
)
"""
_DDL_SUGG = ("CREATE TABLE daily_suggestions (id INTEGER PRIMARY KEY AUTOINCREMENT, "
             "fecha TEXT, tag TEXT, titulo TEXT, cuerpo TEXT, estado TEXT, ref TEXT)")
_DDL_INBOX = ("CREATE TABLE inbox_items (id INTEGER PRIMARY KEY AUTOINCREMENT, "
              "tipo TEXT, texto TEXT, origen TEXT, procesado BOOLEAN)")


def _fake_embed(text_in: str):
    t = (text_in or "").lower()
    return [float(t.count("alfa") + 1), float(t.count("beta") + 1), float(len(t) % 7)]


class _FakeAnalysisClient:
    """Minimal OpenAI-compatible stub: client.chat.completions.create(...).choices[0].message.content"""
    class _Msg:
        content = "Resumen ejecutivo de prueba. Sin cifras inventadas."

    class _Choice:
        message = None

    class _Resp:
        choices = None

    class _Comp:
        def create(self, **kw):
            r = _FakeAnalysisClient._Resp()
            ch = _FakeAnalysisClient._Choice()
            ch.message = _FakeAnalysisClient._Msg()
            r.choices = [ch]
            return r

    class _Chat:
        completions = None

    def __init__(self):
        self.chat = _FakeAnalysisClient._Chat()
        self.chat.completions = _FakeAnalysisClient._Comp()


# ── Extractor: structured office formats ──────────────────────────────────────

def test_extract_xlsx_table_per_sheet():
    openpyxl = pytest.importorskip("openpyxl")
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Caja"
    ws.append(["Concepto", "Monto"])
    ws.append(["Ingresos", "COP 1.000.000"])
    ws2 = wb.create_sheet("Cartera")
    ws2.append(["Cliente", "Saldo"])
    ws2.append(["ACME", "500"])
    buf = io.BytesIO()
    wb.save(buf)
    out = _extract.extract_resource(buf.getvalue(), "libro.xlsx", "")
    assert "### Hoja: Caja" in out
    assert "### Hoja: Cartera" in out
    assert "| Concepto | Monto |" in out
    assert "| Ingresos | COP 1.000.000 |" in out   # amount stays paired with its label
    assert "| ACME | 500 |" in out


def test_extract_pptx_slides_and_notes():
    pytest.importorskip("pptx")
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # "Title Only" has a title placeholder
    slide.shapes.title.text = "Pitch La Falla"
    slide.notes_slide.notes_text_frame.text = "Recordar mencionar el Eje Cafetero"
    buf = io.BytesIO()
    prs.save(buf)
    out = _extract.extract_resource(buf.getvalue(), "deck.pptx", "")
    assert "### Slide 1" in out
    assert "Pitch La Falla" in out
    assert "Notas: Recordar mencionar el Eje Cafetero" in out


def test_extract_csv_as_table():
    data = b"Concepto,Monto\nIngresos,1000\nEgresos,500\n"
    out = _extract.extract_resource(data, "flujo.csv", "text/csv")
    assert "| Concepto | Monto |" in out
    assert "| Ingresos | 1000 |" in out
    assert "| Egresos | 500 |" in out


def test_extract_csv_semicolon_delimiter():
    out = _extract.extract_resource(b"a;b;c\n1;2;3\n", "x.csv", "")
    assert "| a | b | c |" in out
    assert "| 1 | 2 | 3 |" in out


# ── Extractor: images via vision (mocked) ─────────────────────────────────────

def _patch_vision(monkeypatch, content="TEXTO: Factura 123\nDescripción: una factura de proveedor."):
    monkeypatch.setattr(_extract, "VISION_PROVIDER", "openai")
    monkeypatch.setenv("OPENAI_API_KEY", "x")

    class _Msg:
        pass

    class _Choice:
        pass

    class _Resp:
        pass

    class _Comp:
        def create(self, **kw):
            msgs = kw.get("messages", [])
            parts = [p for m in msgs for p in (m.get("content") or []) if isinstance(p, dict)]
            assert any(p.get("type") == "image_url" for p in parts), "vision call must include the image"
            r = _Resp()
            ch = _Choice()
            ch.message = _Msg()
            ch.message.content = content
            r.choices = [ch]
            return r

    class _Chat:
        def __init__(self):
            self.completions = _Comp()

    class _Client:
        def __init__(self, *a, **k):
            self.chat = _Chat()

    import openai
    monkeypatch.setattr(openai, "OpenAI", _Client)


def test_extract_image_vision_mocked(monkeypatch):
    _patch_vision(monkeypatch)
    out = _extract.extract_resource(b"\x89PNG\r\n fake bytes", "foto.png", "image/png")
    assert "Factura 123" in out


def test_extract_image_no_provider_is_honest(monkeypatch):
    monkeypatch.setattr(_extract, "VISION_PROVIDER", "openai")
    monkeypatch.delenv("OPENAI_API_KEY", raising=False)
    assert _extract.extract_resource(b"anything", "foto.jpg", "image/jpeg") == ""


# ── Extractor: links via Firecrawl (mocked) ───────────────────────────────────

class _FakeFirecrawlOK:
    ok = True

    def json(self):
        return {"success": True, "data": {"markdown": "# Politica de datos\nContenido real alfa beta."}}


class _FakeFirecrawlFail:
    ok = False
    status_code = 502
    text = "bad gateway"

    def json(self):
        return {}


def test_extract_url_firecrawl_mocked(monkeypatch):
    import requests
    monkeypatch.setattr(requests, "post", lambda *a, **k: _FakeFirecrawlOK())
    out = _extract.extract_url("https://lafalla.co/politica")
    assert "Contenido real alfa beta." in out


def test_extract_url_failure_is_honest(monkeypatch):
    import requests
    monkeypatch.setattr(requests, "post", lambda *a, **k: _FakeFirecrawlFail())
    assert _extract.extract_url("https://lafalla.co/x") == ""

    def _boom(*a, **k):
        raise RuntimeError("firecrawl down")

    monkeypatch.setattr(requests, "post", _boom)
    assert _extract.extract_url("https://lafalla.co/x") == ""
    assert _extract.extract_url("") == ""


# ── Extractor: honesty + text passthrough + cap ───────────────────────────────

def test_unsupported_binary_is_honest():
    blob = bytes(range(256)) * 4   # ~half are invalid utf-8 -> replacement chars
    assert _extract.extract_resource(blob, "data.bin", "application/octet-stream") == ""


def test_text_passthrough():
    out = _extract.extract_resource("hola alfa mundo".encode(), "n.txt", "text/plain")
    assert "hola alfa mundo" in out


def test_cap_truncates():
    out = _extract.extract_resource(("a" * 100).encode(), "n.txt", "text/plain", max_chars=10)
    assert len(out) == 10


# ── DB state + endpoint (TestClient) ──────────────────────────────────────────

def _setup_db(tmp_path, monkeypatch):
    db = tmp_path / "ing.db"
    url = f"sqlite:///{db.as_posix()}"
    eng = create_engine(url)
    with eng.begin() as c:
        c.execute(text(_DDL_CHUNKS))
        c.execute(text(_DDL_SUGG))
        c.execute(text(_DDL_INBOX))
    monkeypatch.setattr(st, "DATABASE_URL", url)   # handler reads the module global at call time
    monkeypatch.setenv("DATABASE_URL", url)         # rag._db_url reads the env at call time
    monkeypatch.setattr(rag_mod, "get_embedding", _fake_embed)
    monkeypatch.setattr(st, "_analysis_client_model", lambda: (_FakeAnalysisClient(), "gpt-4o-mini"))
    return url, eng


def _client():
    from fastapi.testclient import TestClient
    from api_gateway.main import app
    return TestClient(app)


def test_ingest_resource_xlsx_indexes_chunks(tmp_path, monkeypatch):
    openpyxl = pytest.importorskip("openpyxl")
    url, eng = _setup_db(tmp_path, monkeypatch)
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Presupuesto"
    ws.append(["Concepto", "Monto"])
    ws.append(["Cine", "COP 600.000.000"])
    buf = io.BytesIO()
    wb.save(buf)

    r = _client().post(
        "/api/strategy/ingest-resource",
        headers=H,
        data={"intent": "obs"},
        files={"file": ("presupuesto.xlsx", buf.getvalue(),
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")},
    )
    assert r.status_code == 200, r.text
    assert r.json()["chunks_indexed"] >= 1
    with eng.connect() as c:
        n = c.execute(text("SELECT COUNT(*) FROM document_chunks WHERE source_name='presupuesto.xlsx'")).scalar()
    assert n >= 1


def test_ingest_resource_url_fetches_and_indexes(tmp_path, monkeypatch):
    url, eng = _setup_db(tmp_path, monkeypatch)
    import requests
    monkeypatch.setattr(requests, "post", lambda *a, **k: _FakeFirecrawlOK())

    r = _client().post(
        "/api/strategy/ingest-resource",
        headers=H,
        data={"intent": "obs", "url": "https://lafalla.co/politica"},
    )
    assert r.status_code == 200, r.text
    assert r.json()["chunks_indexed"] >= 1
    with eng.connect() as c:
        n = c.execute(text("SELECT COUNT(*) FROM document_chunks WHERE source_name='https://lafalla.co/politica'")).scalar()
    assert n >= 1


def test_ingest_resource_url_failure_is_honest(tmp_path, monkeypatch):
    url, eng = _setup_db(tmp_path, monkeypatch)
    import requests
    monkeypatch.setattr(requests, "post", lambda *a, **k: _FakeFirecrawlFail())

    r = _client().post(
        "/api/strategy/ingest-resource",
        headers=H,
        data={"intent": "obs", "url": "https://lafalla.co/down"},
    )
    assert r.status_code == 400   # could not read -> honest, no fabricated analysis


def test_ingest_resource_url_requires_key():
    r = _client().post("/api/strategy/ingest-resource", data={"intent": "obs", "url": "https://x.co"})
    assert r.status_code in (401, 403)
