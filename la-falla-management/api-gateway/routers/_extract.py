"""Shared multi-format extractor (change multi-format-ingestion).

One place that turns an uploaded resource (bytes + filename/content-type) or a URL into plain text for
the document RAG and the chat one-shot path. Reused by both upload paths so they never drift again
(design D5): strategy.py `_extract_text` / `ingest-resource`, and chat.py `/extract-file`.

Honesty rule (inherited from wire-document-rag): if a resource cannot be read, return "" so the caller
reports it instead of fabricating content. Never return a placeholder as if it were real content.

Supported: PDF, DOCX, XLSX/XLS, PPTX, CSV, images (via a vision model), and text-like files. URLs are
fetched with Firecrawl (`/v1/scrape`, already running on the VPS).
"""

import io
import os

# Per-resource extraction cap. Structured formats (sheets/decks) can be large; raised from the old 15k.
# Callers may override (chat /extract-file keeps a larger cap for the one-shot read).
MAX_CHARS     = int(os.environ.get("EXTRACT_MAX_CHARS", "30000"))
MAX_PDF_PAGES = int(os.environ.get("EXTRACT_MAX_PDF_PAGES", "50"))

# Images → text via a configured vision model (design D3). Default OpenAI gpt-4o-mini (multimodal).
VISION_PROVIDER = os.environ.get("VISION_PROVIDER", "openai").lower()
VISION_MODEL    = os.environ.get("VISION_MODEL", "gpt-4o-mini")

# Links → Firecrawl on the VPS. Default to the reachable host:published-port (the old 10.0.1.132 is dead
# infra); env can still override. Keeps links working even if FIRECRAWL_URL is not set in the panel.
FIRECRAWL_URL     = os.environ.get("FIRECRAWL_URL", "http://72.61.73.132:3002")
FIRECRAWL_API_KEY = os.environ.get("FIRECRAWL_API_KEY", "")


def _cap(text_in: str, max_chars: int) -> str:
    """Trim to the char cap, logging truncation so there is no silent loss (design D6)."""
    if not text_in:
        return ""
    if len(text_in) > max_chars:
        print(f"[_extract] texto truncado: {len(text_in)} -> {max_chars} chars")
        return text_in[:max_chars]
    return text_in


def _looks_binary(text_in: str) -> bool:
    """Heuristic: a decode dominated by U+FFFD replacement chars is binary, not readable text."""
    if not text_in:
        return False
    bad = text_in.count("�")
    return bad > len(text_in) * 0.3


# ── PDF ─────────────────────────────────────────────────────────────────────────
def _extract_pdf(data: bytes) -> str:
    """pdfplumber (tables row-by-row + prose) → pypdf fallback. Same behavior as wire-document-rag."""
    try:
        import pdfplumber
        parts = []
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            for i, page in enumerate(pdf.pages[:MAX_PDF_PAGES]):
                page_parts = []
                txt = page.extract_text() or ""
                if txt.strip():
                    page_parts.append(txt)
                for table in (page.extract_tables() or []):
                    for row in table:
                        cells = [(c or "").strip() for c in row]
                        if any(cells):
                            page_parts.append(" | ".join(cells))
                if page_parts:
                    parts.append(f"[Pág {i+1}]\n" + "\n".join(page_parts))
        out = "\n\n".join(parts).strip()
        if out:
            return out
    except Exception:
        pass
    try:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(data))
        pages = []
        for i, page in enumerate(reader.pages[:MAX_PDF_PAGES]):
            t = page.extract_text() or ""
            if t.strip():
                pages.append(f"[Pág {i+1}]\n{t}")
        return "\n\n".join(pages).strip()
    except Exception:
        return ""


# ── DOCX ────────────────────────────────────────────────────────────────────────
def _extract_docx(data: bytes) -> str:
    try:
        import zipfile
        import xml.etree.ElementTree as ET
        z = zipfile.ZipFile(io.BytesIO(data))
        if "word/document.xml" in z.namelist():
            xml_data = z.read("word/document.xml").decode("utf-8", errors="ignore")
            root = ET.fromstring(xml_data)
            ns = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
            return " ".join(elem.text for elem in root.iter(f"{ns}t") if elem.text).strip()
    except Exception:
        return ""
    return ""


# ── Excel ───────────────────────────────────────────────────────────────────────
def _extract_xlsx(data: bytes, max_chars: int) -> str:
    """openpyxl read-only/data-only; each sheet linearized row-by-row as a markdown table so each row's
    cells stay associated (amount<->label pairing the budgets need). Stops early at the char cap."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    except Exception:
        return ""
    parts = []
    acc = 0
    try:
        for ws in wb.worksheets:
            rows_out = []
            for row in ws.iter_rows(values_only=True):
                cells = ["" if c is None else str(c).strip() for c in row]
                # drop trailing empties to keep rows tight
                while cells and cells[-1] == "":
                    cells.pop()
                if any(cells):
                    line = "| " + " | ".join(cells) + " |"
                    rows_out.append(line)
                    acc += len(line) + 1
                if acc > max_chars:
                    break
            if rows_out:
                parts.append(f"### Hoja: {ws.title}\n" + "\n".join(rows_out))
            if acc > max_chars:
                break
    finally:
        try:
            wb.close()
        except Exception:
            pass
    return "\n\n".join(parts).strip()


# ── PowerPoint ───────────────────────────────────────────────────────────────────
def _extract_pptx(data: bytes) -> str:
    """python-pptx: per slide, shape text + table cells + speaker notes, in slide order."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
    except Exception:
        return ""
    parts = []
    try:
        for idx, slide in enumerate(prs.slides, start=1):
            chunks = []
            for shape in slide.shapes:
                try:
                    if shape.has_text_frame and shape.text_frame.text.strip():
                        chunks.append(shape.text_frame.text.strip())
                    if shape.has_table:
                        for trow in shape.table.rows:
                            cells = [c.text.strip() for c in trow.cells]
                            if any(cells):
                                chunks.append("| " + " | ".join(cells) + " |")
                except Exception:
                    continue
            try:
                if slide.has_notes_slide:
                    notes = (slide.notes_slide.notes_text_frame.text or "").strip()
                    if notes:
                        chunks.append("Notas: " + notes)
            except Exception:
                pass
            if chunks:
                parts.append(f"### Slide {idx}\n" + "\n".join(chunks))
    except Exception:
        return "\n\n".join(parts).strip()
    return "\n\n".join(parts).strip()


# ── CSV ─────────────────────────────────────────────────────────────────────────
def _extract_csv(data: bytes) -> str:
    """Aligned table rows (sniffed delimiter) instead of a raw blob; honest fallback to raw decode."""
    import csv
    try:
        raw = data.decode("utf-8-sig", errors="replace")
    except Exception:
        return ""
    try:
        sample = raw[:4096]
        try:
            delim = csv.Sniffer().sniff(sample, delimiters=",;\t|").delimiter
        except Exception:
            delim = ","
        out = []
        for row in csv.reader(io.StringIO(raw), delimiter=delim):
            cells = [(c or "").strip() for c in row]
            if any(cells):
                out.append("| " + " | ".join(cells) + " |")
        joined = "\n".join(out).strip()
        return joined or raw.strip()
    except Exception:
        return raw.strip()


# ── Images (vision) ──────────────────────────────────────────────────────────────
def extract_image(data: bytes, content_type: str = "") -> str:
    """Turn an image into indexable text via a configured vision model. No provider/key → "" (honest:
    never fabricate a description). Provider/model via VISION_PROVIDER / VISION_MODEL (design D3)."""
    if VISION_PROVIDER != "openai":
        return ""  # only OpenAI vision wired today; be honest rather than fake it
    key = os.environ.get("OPENAI_API_KEY", "")
    if not key:
        return ""
    import base64
    mime = content_type if (content_type or "").startswith("image/") else "image/png"
    b64 = base64.b64encode(data).decode("ascii")
    data_url = f"data:{mime};base64,{b64}"
    try:
        from openai import OpenAI
        client = OpenAI(api_key=key)
        resp = client.chat.completions.create(
            model=VISION_MODEL,
            max_tokens=900,
            temperature=0.2,
            messages=[{
                "role": "user",
                "content": [
                    {"type": "text", "text": (
                        "Transcribe TODO el texto visible de esta imagen tal cual aparece "
                        "(números, cifras, fechas, etiquetas, encabezados) y luego añade una breve "
                        "descripción estructurada de lo que muestra. No inventes datos que no se vean. "
                        "Responde en español."
                    )},
                    {"type": "image_url", "image_url": {"url": data_url}},
                ],
            }],
        )
        return (resp.choices[0].message.content or "").strip()
    except Exception:
        return ""


# ── URLs (Firecrawl) ─────────────────────────────────────────────────────────────
def extract_url(url: str) -> str:
    """Fetch a URL's readable content as markdown via Firecrawl (real Playwright on the VPS). Failure →
    "" (honest; the caller must not index a placeholder as if it were content). Mirrors chat.py."""
    url = (url or "").strip()
    if not url:
        return ""
    import requests
    headers = {}
    if FIRECRAWL_API_KEY:
        headers["Authorization"] = f"Bearer {FIRECRAWL_API_KEY}"
    try:
        r = requests.post(
            f"{FIRECRAWL_URL}/v1/scrape",
            json={"url": url, "formats": ["markdown"], "onlyMainContent": True},
            headers=headers,
            timeout=30,
        )
        if not r.ok:
            return ""
        data = r.json()
        if not data.get("success"):
            return ""
        content = ((data.get("data", {}) or {}).get("markdown", "") or "").strip()
        return content
    except Exception:
        return ""


# ── Dispatcher ───────────────────────────────────────────────────────────────────
def extract_resource(data: bytes, filename: str = "", content_type: str = "", max_chars: int = None) -> str:
    """Dispatch by extension / content-type to the right extractor; cap the result. Honest "" when the
    format is unsupported or the content cannot be read (no fabrication, no placeholder)."""
    cap = max_chars or MAX_CHARS
    name = (filename or "").lower()
    ctype = (content_type or "").lower()

    if name.endswith(".pdf") or ctype.startswith("application/pdf"):
        return _cap(_extract_pdf(data), cap)

    if name.endswith(".docx") or "wordprocessingml" in ctype:
        return _cap(_extract_docx(data), cap)

    if name.endswith((".xlsx", ".xlsm", ".xls")) or "spreadsheetml" in ctype or ctype == "application/vnd.ms-excel":
        return _cap(_extract_xlsx(data, cap), cap)

    if name.endswith(".pptx") or "presentationml" in ctype:
        return _cap(_extract_pptx(data), cap)

    if name.endswith(".csv") or ctype == "text/csv":
        return _cap(_extract_csv(data), cap)

    if name.endswith((".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp")) or ctype.startswith("image/"):
        return _cap(extract_image(data, content_type), cap)

    if ctype.startswith("text/") or name.endswith(
        (".txt", ".md", ".markdown", ".json", ".sql", ".py", ".js", ".ts", ".html", ".htm", ".xml", ".yaml", ".yml", ".csv")
    ):
        try:
            return _cap(data.decode("utf-8", errors="replace"), cap)
        except Exception:
            return ""

    # Unknown type: try utf-8, but stay honest about binary blobs.
    try:
        decoded = data.decode("utf-8", errors="replace")
    except Exception:
        return ""
    if _looks_binary(decoded):
        return ""
    return _cap(decoded, cap)
